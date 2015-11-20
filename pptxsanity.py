#!/usr/bin/env python
# -*- coding: utf-8 -*-
#
# With code by Eric Jang ericjang2004@gmail.com
TIMEOUT=6 # URL request timeout in seconds
MAXRETRY=4
MAXREDIR=4
SKIP200=1

from pptx import Presentation
import sys
import re
import os
import shutil
import glob
import tempfile
import urllib3
try:
    import urllib3.contrib.pyopenssl
    urllib3.contrib.pyopenssl.inject_into_urllib3()
except ImportError:
    pass

import signal
from zipfile import ZipFile
from xml.dom.minidom import parse
import xml.etree.ElementTree as etree
import platform
import ssl
from functools import wraps

def sslwrap(func):
    @wraps(func)
    def bar(*args, **kw):
        kw['ssl_version'] = ssl.PROTOCOL_TLSv1
        return func(*args, **kw)
    return bar



# Remove trailing unwanted characters from the end of URL's
# This is a recursive function. Did I do it well? I don't know.
def striptrailingchar(s):
    # The valid URL charset is A-Za-z0-9-._~:/?#[]@!$&'()*+,;= and & followed by hex character
    # I don't have a better way to parse URL's from the cruft that I get from XML content, so I
    # also remove .),;'? too.  Note that this is only the end of the URL (making ? OK to remove)
    #if s[-1] not in "ABCDEFGHIJKLMNOPQRSTUVWXYZZabcdefghijklmnopqrstuvwxyzz0123456789-_~:/#[]@!$&(*+=":
    if s[-1] not in "ABCDEFGHIJKLMNOPQRSTUVWXYZZabcdefghijklmnopqrstuvwxyzz0123456789-_~:#[]@!$&(*+=":
        s = striptrailingchar(s[0:-1])
    elif s[-5:] == "&quot":
        s = striptrailingchar(s[0:-5])
    else:
        pass
    return s

def parseslidenotes(pptxfile):
    urls = []
    tmpd = tempfile.mkdtemp()

    ZipFile(pptxfile).extractall(path=tmpd, pwd=None)
    path = tmpd + '/ppt/notesSlides/'

    for infile in glob.glob(os.path.join(path, '*.xml')):
        #parse each XML notes file from the notes folder.

        # Get the slide number
        dom = parse(infile)
        noteslist = dom.getElementsByTagName('a:t')
        slideNumber = noteslist.pop()
        slideNumber = slideNumber.toxml().replace('<a:t>', '').replace('</a:t>', '')

        # Parse slide notes, adding a space after each paragraph marker, and removing XML markup
        s=dom.toxml().encode('ascii', 'ignore').replace("</a:p>","</a:p> ")
        text = ''.join(etree.fromstring(s).itertext())

        # Parse URL content from notes text
        urlmatches = re.findall(urlmatchre,text)
        if len(urlmatches) > 0:
            for match in urlmatches: # Now it's a tuple
                 for urlmatch in match:
                      if urlmatch != '':
                          urls.append(striptrailingchar(urlmatch))

    # Remove all the files created with unzip
    shutil.rmtree(tmpd)
    return urls

# Parse the text on slides using the python-pptx module, return URLs
def parseslidetext(prs):
    urls = []
    nexttitle = False
    for slide in prs.slides:
        text_runs = []
        for shape in slide.shapes:
            try:
                if not shape.has_text_frame:
                    continue
            except AttributeError:
                sys.stderr.write("Error: Please upgrade your version of python-pptx: pip uninstall python-pptx ; pip install python-pptx\n")
                sys.exit(-1)
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)

            for text in text_runs:
                if text == None : continue
                try:
                    m = re.match(urlmatchre,text)
                except IndexError,TypeError:
                    continue
                if m != None:
                    url = striptrailingchar(m.groups()[0])
                    if url not in urls:
                        urls.append(url)
    return urls

def signal_exit(signal, frame):
    sys.exit(0)

if __name__ == "__main__":
    if (len(sys.argv) != 2):
        print "Validate URLs in the notes and slides of a PowerPoint pptx file."
        print "Check GitHub for updates: http://github.com/joswr1ght/pptxsanity\n"
        if (platform.system() == 'Windows'):
            print "Usage: pptxsanity.exe [pptx file]"
        else:
            print "Usage: pptxsanity.py [pptx file]"
        sys.exit(1)

    signal.signal(signal.SIGINT, signal_exit)

    # Disable urllib3 InsecureRequestWarning
    try:
        urllib3.disable_warnings()
    except AttributeError:
        sys.stdout.write("You need to upgrade your version of the urllib3 library to the latest available.\n");
        sys.stdout.write("Try running the following command to upgrade urllib3: sudo pip install urllib3 --upgrade\n");
        sys.exit(1)
    
    try:
        prs = Presentation(sys.argv[1])
    except Exception:
        sys.stderr.write("Invalid PPTX file: " + sys.argv[1] + "\n")
        sys.exit(-1)
    
    # This may be the most insane regex I've ever seen.  It's very comprehensive, but it's too aggressive for
    # what I want.  It matches arp:remote in ettercap -TqM arp:remote // //, so I'm using something simpler
    #urlmatchre = re.compile(r"""((?:[a-z][\w-]+:(?:/{1,3}|[a-z0-9%])|www\d{0,3}[.]|[a-z0-9.\-]+[.‌​][a-z]{2,4}/)(?:[^\s()<>]+|(([^\s()<>]+|(([^\s()<>]+)))*))+(?:(([^\s()<>]+|(‌​([^\s()<>]+)))*)|[^\s`!()[]{};:'".,<>?«»“”‘’]))""", re.DOTALL)
    urlmatchre = re.compile(r'((https?://[^\s<>"]+|www\.[^\s<>"]+))',re.DOTALL)
    privateaddr = re.compile(r'(\S+127\.)|(\S+192\.168\.)|(\S+10\.)|(\S+172\.1[6-9]\.)|(\S+172\.2[0-9]\.)|(\S+172\.3[0-1]\.)|(\S+::1)')

    SKIP200=int(os.getenv('SKIP200', 1))
    
    urls = []
    urls += parseslidetext(prs)
    urls += parseslidenotes(sys.argv[1])

    # De-duplicate URL's
    urls = list(set(urls))

    for url in urls:
        if "whois.net" in url:
            continue

        url = url.encode('ascii', 'ignore')

        # Add default URI for www.anything
        if url[0:3] == "www": url="http://"+url

        # Skip private IP addresses
        if re.match(privateaddr,url): continue

        headers = { 'User-Agent' : 'Mozilla/5.0 (Windows NT 6.1; Win64; x64; rv:35.0) Gecko/20100101 Firefox/35.0' }
        retries=urllib3.Retry(redirect=4, total=4, connect=0, read=0)
        http = urllib3.PoolManager(timeout=6, retries=retries)
        try:
            req=http.request('HEAD', url, headers=headers)
            code=req.status
        except Exception, e:
            print "ERR : " + url
            continue

        # Some websites return 404 for HEAD requests (microsoft.com).  If we get a 404, try to retrieve using GET
        # and report the corresponding response code.  Also check out 405 "Method not allowed" responses.
        if code == 404 or code == 405:
            # Stupid non-compliant web server
            try:
                req=http.request('GET', url, headers=headers)
                code=req.status
            except Exception, e:
                print "ERR : " + url
                continue

        if code == 200 and SKIP200 == 1:
            continue
        print str(code) + " : " + url
